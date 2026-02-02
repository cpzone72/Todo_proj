from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Define content structure
    slides_content = [
        {
            "title": "NVR3864-V6-IQ 신규 도입 BMT 테스트 계획",
            "content": "성능, 안정성, AI 기능 검증을 위한 벤치마킹 테스트 절차\n\n대상 장비: NVR3864-V6-IQ (64ch, 8HDD, 2U)",
            "layout": 0  # Title Slide
        },
        {
            "title": "1. BMT 개요 및 목적",
            "content": "■ 대상 장비\n- 모델명: NVR3864-V6-IQ\n- 주요 스펙: 64채널 입력, 8 SATA HDD, RAID 지원, 4K 출력\n\n■ 목적\n- 64채널 대용량 데이터 처리 능력 검증\n- AI 분석 기능(VCA) 정확도 확인\n- RAID 및 Failover 안정성 확보를 통한 현장 도입 적합성 판단",
            "layout": 1  # Title and Content
        },
        {
            "title": "2. 테스트 환경 구성 (권장)",
            "content": "■ NVR\n- NVR3864-V6-IQ 본체 (HDD 8TB x 8EA, RAID 5/6 권장)\n\n■ 카메라 및 입력 소스\n- 실물 카메라: 4K(8MP) AI 카메라 4대 이상 (지능형 기능 검증용)\n- 시뮬레이터: RTSP 64채널 풀 로드 (H.265, 4Mbps 기준)\n\n■ 디스플레이 및 네트워크\n- 4K 모니터 2대 (HDMI 1/2 독립 출력)\n- L2 기가비트 스위치, 클라이언트 PC 1대",
            "layout": 1
        },
        {
            "title": "3-A. 기본 성능 및 UI (General Performance)",
            "content": "A-1. 부팅 속도\n   - 콜드 부팅 후 녹화 시작까지 3분 이내\n\nA-2. 출력 해상도\n   - HDMI 1/2 독립 4K 출력 확인 (60Hz/30Hz)\n\nA-3. UI 반응성\n   - 64채널 라이브 상태에서 메뉴 진입/전환 시 지연(Lag) 없음 확인",
            "layout": 1
        },
        {
            "title": "3-B. 영상 처리 및 부하 (Processing & Load)",
            "content": "B-1. 입력 대역폭 한계 테스트\n   - 384Mbps(Smart Off) 부하 인가 시 프레임 드랍 여부\n\nB-2. 압축 효율 검증\n   - Ultra 265 vs H.264 비트레이트 절감률 비교 (약 50%)\n\nB-3. 디코딩 능력\n   - 4K 8채널 또는 2MP 32채널 동시 라이브 시청 시 부드러움 확인\n\nB-4. 장기 부하 테스트\n   - 72시간 연속 가동 후 시스템 로그(재부팅, 에러) 점검",
            "layout": 1
        },
        {
            "title": "3-C. 저장 및 안정성 (Storage & Reliability)",
            "content": "C-1. RAID 구성\n   - RAID 5/6 볼륨 생성 및 설정 용이성\n\nC-2. HDD 장애 대응 (Rebuild)\n   - 녹화 중 HDD 강제 탈거 후 리빌딩 시 녹화 지속 여부 확인\n\nC-3. Hot Spare (N+1)\n   - 예비 디스크/장비 자동 전환 테스트\n\nC-4. ANR (네트워크 단절 보정)\n   - 네트워크 재연결 시 카메라 SD카드 데이터 자동 백업 확인",
            "layout": 1
        },
        {
            "title": "3-D. 지능형 기능 및 검색 (AI & Search)",
            "content": "D-1. VCA 분석 (침입/라인 크로싱)\n   - 구역 설정 후 객체 진입 시 알람 정확도 (90% 이상)\n\nD-2. 객체 구분 (Smart Intrusion Prevention)\n   - 사람/차량 분류 필터링 성능 (동물, 나뭇잎 오작동 제외)\n\nD-3. 스마트 검색 (AcuSearch)\n   - 속성 검색(빨간 옷, 차량 등) 시 검색 속도 측정\n\nD-4. 얼굴 인식 (Face Detection)\n   - 등록된 얼굴 매칭 및 즉시 알람 표출 여부",
            "layout": 1
        },
        {
            "title": "3-E. 네트워크 및 보안 / 3-F. 호환성",
            "content": "■ 네트워크 및 보안\n- E-1. 원격 접속: 모바일 앱 16분할 라이브 로딩 속도 (3초 이내)\n- E-2. 이중화 네트워크: 듀얼 LAN 망 분리 또는 로드 밸런싱\n- E-3. 보안: HTTPS 및 802.1x 적용 시 영상 전송 안정성\n\n■ 호환성\n- F-1. 타사 카메라 ONVIF 연동 (라이브, 녹화, PTZ 제어 확인)",
            "layout": 1
        },
        {
            "title": "4. BMT 결과 보고 및 팁",
            "content": "■ 결과 보고서 항목\n- 종합 판정 (적합 / 부적합)\n- 주요 이슈 사항 및 개선 필요점\n- 운영 권고 사항 (HDD 등급, 네트워크 구성 등)\n\n■ 테스트 팁\n- 대량 설정 시 EZTools/Guard Station CMS 활용 권장\n- 'Smart Mode' On/Off에 따른 대역폭 차이(200M vs 384M) 비교 필수",
            "layout": 1
        }
    ]

    for slide_info in slides_content:
        layout = prs.slide_layouts[slide_info["layout"]]
        slide = prs.slides.add_slide(layout)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_info["title"]
        
        # Set content
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = slide_info["content"]

    file_name = "NVR3864-V6-IQ_BMT_Test_Plan.pptx"
    prs.save(file_name)
    print(f"Presentation saved as {file_name}")

if __name__ == "__main__":
    create_presentation()
